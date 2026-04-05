"""
PPT & PDF Analyzer Module for EchoMind
Extracts topics from PPTX/PDF presentations and generates YouTube playlist recommendations.
"""
import json


def extract_pdf_text(uploaded_file):
    """Extract text from a PDF file, page by page."""
    from io import BytesIO
    
    if hasattr(uploaded_file, 'read'):
        uploaded_file.seek(0)
        file_bytes = uploaded_file.read()
    else:
        file_bytes = uploaded_file
    
    pages_data = []
    
    # Try PyPDF2 first
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(BytesIO(file_bytes))
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages_data.append({"slide": page_num, "content": text.strip()})
        if pages_data:
            return pages_data
    except ImportError:
        pass
    except Exception:
        pass
    
    # Fallback: try pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages_data.append({"slide": page_num, "content": text.strip()})
        if pages_data:
            return pages_data
    except ImportError:
        pass
    except Exception:
        pass
    
    # If nothing worked
    if not pages_data:
        raise ImportError("PDF extraction requires PyPDF2 or pdfplumber. Install: pip install PyPDF2")
    return pages_data

def extract_ppt_text(uploaded_file):
    """Extract text from all slides of a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required. Install it: pip install python-pptx")
    
    from io import BytesIO
    
    if hasattr(uploaded_file, 'read'):
        uploaded_file.seek(0)
        file_bytes = uploaded_file.read()
        prs = Presentation(BytesIO(file_bytes))
    else:
        prs = Presentation(BytesIO(uploaded_file))
    
    slides_data = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides_data.append({
                "slide": slide_num,
                "content": "\n".join(slide_text)
            })
    
    return slides_data


def analyze_ppt_topics(slides_data, groq_client, model="llama-3.3-70b-versatile"):
    """
    Use LLM to analyze PPT content and identify:
    - Topics covered
    - Teaching method used per topic
    - Coverage depth (brief/moderate/comprehensive)
    """
    all_text = ""
    for slide in slides_data:
        all_text += f"\n--- Slide {slide['slide']} ---\n{slide['content']}\n"
    
    if not all_text.strip():
        return {"topics": [], "summary": "No readable text found in the presentation."}
    
    completion = groq_client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "system",
                "content": """You are an educational content analyst. Analyze the presentation slides and return a JSON object with:
{
  "course_title": "Best guess at the course/subject name",
  "summary": "2-3 sentence summary of what the PPT covers",
  "topics": [
    {
      "name": "Topic name",
      "slides": "e.g. 1-5",
      "teaching_method": "e.g. Lecture with diagrams, Case study, Code walkthrough, Theory + examples, Problem solving",
      "coverage": "brief | moderate | comprehensive",
      "key_concepts": ["concept1", "concept2"]
    }
  ]
}

Be thorough. Identify ALL distinct topics."""
            },
            {
                "role": "user",
                "content": f"Analyze this presentation:\n\n{all_text[:8000]}"
            }
        ],
        temperature=0.3,
        max_tokens=4000,
        response_format={"type": "json_object"},
    )
    
    try:
        return json.loads(completion.choices[0].message.content)
    except json.JSONDecodeError:
        return {"topics": [], "summary": "Could not parse analysis results."}


def generate_playlist(analysis, groq_client, model="llama-3.3-70b-versatile"):
    """
    Generate YouTube search queries and playlist recommendations for each topic.
    """
    topics_text = ""
    for topic in analysis.get("topics", []):
        topics_text += f"- {topic['name']} (coverage: {topic.get('coverage', 'unknown')}, method: {topic.get('teaching_method', 'unknown')})\n"
    
    if not topics_text:
        return {"playlist": []}
    
    course_title = analysis.get("course_title", "Unknown Course")
    
    completion = groq_client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "system",
                "content": """You are a study resource curator. Given course topics from a PPT, generate YouTube search queries that will find the best educational videos for each topic.

Return JSON:
{
  "playlist": [
    {
      "topic": "Topic name from PPT",
      "search_query": "Exact YouTube search query to find a great explanation",
      "why": "Brief reason why this video type is recommended",
      "watch_order": 1
    }
  ]
}

Order by logical learning sequence. Make search queries specific and educational (include keywords like "explained", "tutorial", "lecture")."""
            },
            {
                "role": "user",
                "content": f"Course: {course_title}\n\nTopics to cover:\n{topics_text}"
            }
        ],
        temperature=0.3,
        max_tokens=4000,
        response_format={"type": "json_object"},
    )
    
    try:
        return json.loads(completion.choices[0].message.content)
    except json.JSONDecodeError:
        return {"playlist": []}


def render_ppt_analysis(analysis, playlist_data):
    """Render the PPT analysis and playlist in Streamlit."""
    import streamlit as st
    
    if not analysis:
        st.warning("No analysis available.")
        return
    
    # Course overview
    course_title = analysis.get("course_title", "Unknown Course")
    st.markdown(f"#### 📚 {course_title}")
    st.markdown(analysis.get("summary", ""))
    
    # Topics table
    topics = analysis.get("topics", [])
    if topics:
        st.markdown("---")
        st.markdown("#### 📋 Topics Identified")
        for i, topic in enumerate(topics, 1):
            coverage_emoji = {"brief": "🟡", "moderate": "🟠", "comprehensive": "🟢"}.get(topic.get("coverage", ""), "⚪")
            with st.expander(f"{i}. {topic['name']} {coverage_emoji}"):
                st.markdown(f"**Teaching Method:** {topic.get('teaching_method', 'N/A')}")
                st.markdown(f"**Coverage:** {topic.get('coverage', 'N/A').title()}")
                st.markdown(f"**Slides:** {topic.get('slides', 'N/A')}")
                concepts = topic.get("key_concepts", [])
                if concepts:
                    st.markdown("**Key Concepts:** " + ", ".join(concepts))
    
    # Playlist
    playlist = playlist_data.get("playlist", [])
    if playlist:
        st.markdown("---")
        st.markdown("#### 🎬 Recommended YouTube Playlist")
        st.markdown("*Watch in this order for the best learning experience:*")
        
        for item in sorted(playlist, key=lambda x: x.get("watch_order", 99)):
            order = item.get("watch_order", "?")
            topic = item.get("topic", "Unknown")
            query = item.get("search_query", "")
            why = item.get("why", "")
            
            yt_url = f"https://www.youtube.com/results?search_query={query.replace(' ', '+')}"
            
            st.markdown(f"""
**{order}. {topic}**  
🔍 [Search: *{query}*]({yt_url})  
💡 {why}
""")
